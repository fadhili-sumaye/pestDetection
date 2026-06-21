import os
import sys

def check_and_install_dependencies():
    """Checks for python-pptx and installs it if missing."""
    try:
        import pptx
        print("python-pptx is already installed.")
    except ImportError:
        print("python-pptx is not installed. Installing it now via pip...")
        import subprocess
        try:
            subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx"], check=True)
            print("Successfully installed python-pptx!")
        except Exception as e:
            print(f"Error installing python-pptx: {e}")
            print("Please run: pip install python-pptx")
            sys.exit(1)

check_and_install_dependencies()

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # 1. Initialize Presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Blank layout is usually layout index 6 in standard templates
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette
    COLOR_PRIMARY_GREEN = RGBColor(27, 94, 32)      # Deep Forest Green (#1B5E20)
    COLOR_MINT_LIGHT = RGBColor(232, 245, 233)     # Light Mint (#E8F5E9)
    COLOR_DARK_TEXT = RGBColor(33, 33, 33)         # Slate Black (#212121)
    COLOR_WHITE = RGBColor(255, 255, 255)          # White (#FFFFFF)
    COLOR_ACCENT_GOLD = RGBColor(245, 127, 23)     # Accent Gold (#F57F17)
    COLOR_SOFT_GRAY = RGBColor(245, 245, 245)      # Light Gray (#F5F5F5)
    COLOR_MUTED_GREEN = RGBColor(70, 130, 80)      # Muted Green for tags

    # Font names
    FONT_TITLE = "Arial"
    FONT_BODY = "Arial"

    # --- Helper Functions ---
    def set_solid_background(slide, color):
        """Sets a full slide solid background using a rectangle."""
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background() # No border
        # Send shape to back using internal helper
        slide.shapes._spTree.remove(rect._element)
        slide.shapes._spTree.insert(2, rect._element)
        return rect

    def add_header(slide, title_text, dark_mode=False):
        """Adds a standard header to a slide."""
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_TITLE
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE if dark_mode else COLOR_PRIMARY_GREEN

    def add_card(slide, left, top, width, height, fill_color, border_color=None):
        """Creates a background card for content grouping."""
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = fill_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.fill.background()
        return card

    # ==========================================
    # SLIDE 1: Title Slide (Dark Theme)
    # ==========================================
    slide = prs.slides.add_slide(blank_slide_layout)
    set_solid_background(slide, COLOR_PRIMARY_GREEN)

    # Accent decorative shape
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.2), Inches(1.5), Inches(0.1))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT_GOLD
    accent_bar.line.fill.background()

    # Title & Subtitle text frame
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(11.833), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "SMART PEST DETECTION SYSTEM"
    p.font.name = FONT_TITLE
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(20)

    p2 = tf.add_paragraph()
    p2.text = "Real-Time Mobile & AI Agricultural Diagnostics for Healthy Crops"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_MINT_LIGHT
    p2.space_after = Pt(40)

    p3 = tf.add_paragraph()
    p3.text = "Prepared by: Mobile & AI Development Team"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_ACCENT_GOLD

    # ==========================================
    # SLIDE 2: The Challenge (Problem)
    # ==========================================
    slide = prs.slides.add_slide(blank_slide_layout)
    add_header(slide, "The Challenge: Pest Devastation")

    # Left Column: The Problem Points
    txBox_left = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left = txBox_left.text_frame
    tf_left.word_wrap = True

    problems = [
        ("Delayed Diagnostics", "Farmers notice pest infestations only after significant damage has spread across crops."),
        ("Expert Deficit", "Agricultural extension officers are rarely available for immediate, real-time farm inspections."),
        ("Indiscriminate Spraying", "Incorrect diagnosis leads farmers to use wrong pesticides, increasing costs and damaging local ecosystems.")
    ]

    for title, desc in problems:
        p_title = tf_left.add_paragraph() if tf_left.paragraphs[0].text else tf_left.paragraphs[0]
        p_title.text = "• " + title
        p_title.font.name = FONT_TITLE
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY_GREEN
        p_title.space_after = Pt(4)

        p_desc = tf_left.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(15)
        p_desc.font.color.rgb = COLOR_DARK_TEXT
        p_desc.space_after = Pt(20)

    # Right Column: Big Stat Card
    add_card(slide, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.5), COLOR_MINT_LIGHT)
    
    txBox_right = slide.shapes.add_textbox(Inches(7.8), Inches(2.2), Inches(4.4), Inches(3.7))
    tf_right = txBox_right.text_frame
    tf_right.word_wrap = True

    p_stat_num = tf_right.paragraphs[0]
    p_stat_num.text = "20% - 40%"
    p_stat_num.font.name = FONT_TITLE
    p_stat_num.font.size = Pt(56)
    p_stat_num.font.bold = True
    p_stat_num.font.color.rgb = COLOR_ACCENT_GOLD
    p_stat_num.alignment = PP_ALIGN.CENTER
    p_stat_num.space_after = Pt(10)

    p_stat_lbl = tf_right.add_paragraph()
    p_stat_lbl.text = "Global Crop Loss Annually"
    p_stat_lbl.font.name = FONT_TITLE
    p_stat_lbl.font.size = Pt(20)
    p_stat_lbl.font.bold = True
    p_stat_lbl.font.color.rgb = COLOR_PRIMARY_GREEN
    p_stat_lbl.alignment = PP_ALIGN.CENTER
    p_stat_lbl.space_after = Pt(15)

    p_stat_desc = tf_right.add_paragraph()
    p_stat_desc.text = "According to the FAO, pests destroy billions of dollars in cereal crops annually. Early detection is critical to halting localized outbreaks."
    p_stat_desc.font.name = FONT_BODY
    p_stat_desc.font.size = Pt(14)
    p_stat_desc.font.color.rgb = COLOR_DARK_TEXT
    p_stat_desc.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 3: Solution Overview
    # ==========================================
    slide = prs.slides.add_slide(blank_slide_layout)
    add_header(slide, "The Solution: Simple & Automated Diagnostics")

    steps = [
        ("1. CAPTURE", "Capture photo with camera or choose from Android gallery.", "Camera Integration"),
        ("2. ANALYZE", "OkHttp sends image to custom YOLOv8 AI server API.", "Fast Inference"),
        ("3. RESOLVE", "Instantly read pest type, confidence, and target treatment.", "Knowledge Base")
    ]

    card_width = Inches(3.6)
    card_height = Inches(4.5)
    gap = Inches(0.5)
    start_left = Inches(0.75)
    top_pos = Inches(1.8)

    for i, (title, desc, tag) in enumerate(steps):
        left_pos = start_left + i * (card_width + gap)
        # Background card
        add_card(slide, left_pos, top_pos, card_width, card_height, COLOR_SOFT_GRAY)
        
        # Top banner for card
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, top_pos, card_width, Inches(0.8))
        banner.fill.solid()
        banner.fill.fore_color.rgb = COLOR_PRIMARY_GREEN
        banner.line.fill.background()
        
        # Text inside top banner
        txBox_banner = slide.shapes.add_textbox(left_pos, top_pos, card_width, Inches(0.8))
        p_b = txBox_banner.text_frame.paragraphs[0]
        p_b.text = title
        p_b.font.name = FONT_TITLE
        p_b.font.size = Pt(18)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_WHITE
        p_b.alignment = PP_ALIGN.CENTER
        
        # Body text inside card
        txBox_body = slide.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(1.2), card_width - Inches(0.4), Inches(2.8))
        tf_body = txBox_body.text_frame
        tf_body.word_wrap = True
        
        p_btext = tf_body.paragraphs[0]
        p_btext.text = desc
        p_btext.font.name = FONT_BODY
        p_btext.font.size = Pt(16)
        p_btext.font.color.rgb = COLOR_DARK_TEXT
        p_btext.space_after = Pt(30)
        
        p_tag = tf_body.add_paragraph()
        p_tag.text = f"Key: {tag}"
        p_tag.font.name = FONT_BODY
        p_tag.font.bold = True
        p_tag.font.size = Pt(14)
        p_tag.font.color.rgb = COLOR_ACCENT_GOLD

    # ==========================================
    # SLIDE 4: The Custom YOLOv8 Model
    # ==========================================
    slide = prs.slides.add_slide(blank_slide_layout)
    add_header(slide, "AI Core: Custom YOLOv8 Object Detection")

    # Left Box - Features
    add_card(slide, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5), COLOR_MINT_LIGHT)
    tx_model_l = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.1), Inches(4.1))
    tf_ml = tx_model_l.text_frame
    tf_ml.word_wrap = True

    p_ml_title = tf_ml.paragraphs[0]
    p_ml_title.text = "Model Architecture"
    p_ml_title.font.name = FONT_TITLE
    p_ml_title.font.size = Pt(22)
    p_ml_title.font.bold = True
    p_ml_title.font.color.rgb = COLOR_PRIMARY_GREEN
    p_ml_title.space_after = Pt(15)

    model_points = [
        ("Base Architecture:", "YOLOv8 Nano (yolov8n.pt) for rapid millisecond server inference."),
        ("Target Pests:", "Specifically trained to detect destructive cereal crop pests: planthoppers, rice bugs, leaf folders, whorl maggots, and stem borers."),
        ("Image Processing:", "Input is dynamically normalized to 640x640 resolution, ensuring precise localization and boundary identification.")
    ]

    for lbl, text in model_points:
        p = tf_ml.add_paragraph()
        p.text = f"• {lbl} "
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK_TEXT
        
        run = p.add_run()
        run.text = text
        run.font.bold = False
        p.space_after = Pt(10)

    # Right Box - Training Setup
    add_card(slide, Inches(6.983), Inches(1.8), Inches(5.6), Inches(4.5), COLOR_SOFT_GRAY)
    tx_model_r = slide.shapes.add_textbox(Inches(7.233), Inches(2.0), Inches(5.1), Inches(4.1))
    tf_mr = tx_model_r.text_frame
    tf_mr.word_wrap = True

    p_mr_title = tf_mr.paragraphs[0]
    p_mr_title.text = "Training Specifications"
    p_mr_title.font.name = FONT_TITLE
    p_mr_title.font.size = Pt(22)
    p_mr_title.font.bold = True
    p_mr_title.font.color.rgb = COLOR_PRIMARY_GREEN
    p_mr_title.space_after = Pt(15)

    training_points = [
        ("Framework Stack", "PyTorch combined with the Ultralytics SDK"),
        ("Transfer Learning", "Pre-trained MS COCO weights adjusted for agricultural pest features"),
        ("Dataset Structure", "Split into Train, Validation, and Test directories with Roboflow YAML configuration"),
        ("Output Generation", "Compiles into best.pt PyTorch model weights used by our Flask server API")
    ]

    for lbl, text in training_points:
        p = tf_mr.add_paragraph()
        p.text = f"✔  {lbl}\n"
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT_GOLD
        
        run = p.add_run()
        run.text = f"    {text}"
        run.font.bold = False
        run.font.color.rgb = COLOR_DARK_TEXT
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 5: The Android Mobile App
    # ==========================================
    slide = prs.slides.add_slide(blank_slide_layout)
    add_header(slide, "The Android App: Powerful Client UI")

    # Left Column: Core mobile features
    txBox_app_l = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_app_l = txBox_app_l.text_frame
    tf_app_l.word_wrap = True

    mobile_features = [
        ("Camera & Gallery Integration", "Uses modern Android Activity Result APIs (TakePicture & StartActivityForResult) avoiding deprecated callbacks. Managed securely via FileProvider to guarantee camera permissions."),
        ("OkHttp Network Engine", "Handles HTTP multi-part requests. Configured with specific 30-second timeouts to support robust large-image upload over unstable or slow farm connections."),
        ("Dynamic UI Rendering", "Processes the server's structured JSON payload to display the primary pest detected, confidence score (%), and custom treatments without lagging.")
    ]

    for head, text in mobile_features:
        p_h = tf_app_l.add_paragraph() if tf_app_l.paragraphs[0].text else tf_app_l.paragraphs[0]
        p_h.text = f"✔  {head}"
        p_h.font.name = FONT_TITLE
        p_h.font.size = Pt(18)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_PRIMARY_GREEN
        p_h.space_after = Pt(4)

        p_t = tf_app_l.add_paragraph()
        p_t.text = text
        p_t.font.name = FONT_BODY
        p_t.font.size = Pt(14)
        p_t.font.color.rgb = COLOR_DARK_TEXT
        p_t.space_after = Pt(16)

    # Right Column: Visual App Mockup Card
    add_card(slide, Inches(7.5), Inches(1.8), Inches(5.083), Inches(4.5), COLOR_PRIMARY_GREEN)
    
    txBox_app_mock = slide.shapes.add_textbox(Inches(7.8), Inches(2.2), Inches(4.483), Inches(3.7))
    tf_mock = txBox_app_mock.text_frame
    tf_mock.word_wrap = True

    mock_lines = [
        ("Pest Scanner App", 24, COLOR_ACCENT_GOLD, True),
        ("-----------------------", 16, COLOR_MINT_LIGHT, False),
        ("Status: Connected to API", 14, COLOR_WHITE, False),
        ("\n[ Image Selected ]", 16, COLOR_MINT_LIGHT, True),
        ("Filename: image.jpg", 12, COLOR_WHITE, False),
        ("\nResult: Brown Planthopper", 18, COLOR_WHITE, True),
        ("Confidence: 94%", 16, COLOR_ACCENT_GOLD, True),
        ("\nAdvice: Avoid excess Nitrogen, apply Buprofezin.", 13, COLOR_MINT_LIGHT, False)
    ]

    for text, size, color, is_bold in mock_lines:
        p = tf_mock.add_paragraph() if tf_mock.paragraphs[0].text else tf_mock.paragraphs[0]
        p.text = text
        p.font.name = FONT_BODY
        p.font.size = Pt(size)
        p.font.bold = is_bold
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 6: System Architecture & Data Flow
    # ==========================================
    slide = prs.slides.add_slide(blank_slide_layout)
    add_header(slide, "System Architecture & Data Flow")

    # Step 1 box
    add_card(slide, Inches(0.75), Inches(2.5), Inches(2.5), Inches(2.5), COLOR_SOFT_GRAY)
    tx1 = slide.shapes.add_textbox(Inches(0.85), Inches(2.7), Inches(2.3), Inches(2.1))
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    tf1.paragraphs[0].text = "1. Mobile Client\n\n- Android App\n- Capture Photo\n- FileProvider URI"
    tf1.paragraphs[0].font.size = Pt(14)
    tf1.paragraphs[0].font.color.rgb = COLOR_DARK_TEXT
    tf1.paragraphs[0].font.bold = True

    # Arrow 1
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.4), Inches(3.4), Inches(0.6), Inches(0.4))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = COLOR_ACCENT_GOLD
    arrow1.line.fill.background()

    # Step 2 box
    add_card(slide, Inches(4.15), Inches(2.5), Inches(2.5), Inches(2.5), COLOR_MINT_LIGHT)
    tx2 = slide.shapes.add_textbox(Inches(4.25), Inches(2.7), Inches(2.3), Inches(2.1))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = "2. Network Upload\n\n- OkHttp Post\n- Multipart Body\n- Server Port 5000"
    tf2.paragraphs[0].font.size = Pt(14)
    tf2.paragraphs[0].font.color.rgb = COLOR_PRIMARY_GREEN
    tf2.paragraphs[0].font.bold = True

    # Arrow 2
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.8), Inches(3.4), Inches(0.6), Inches(0.4))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = COLOR_ACCENT_GOLD
    arrow2.line.fill.background()

    # Step 3 box
    add_card(slide, Inches(7.55), Inches(2.5), Inches(2.5), Inches(2.5), COLOR_PRIMARY_GREEN)
    tx3 = slide.shapes.add_textbox(Inches(7.65), Inches(2.7), Inches(2.3), Inches(2.1))
    tf3 = tx3.text_frame
    tf3.word_wrap = True
    tf3.paragraphs[0].text = "3. AI Backend\n\n- Flask Server\n- YOLOv8 Inference\n- JSON Response"
    tf3.paragraphs[0].font.size = Pt(14)
    tf3.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf3.paragraphs[0].font.bold = True

    # Arrow 3 (Return flow indicated)
    arrow3 = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(4.15), Inches(5.5), Inches(5.0), Inches(0.4))
    arrow3.fill.solid()
    arrow3.fill.fore_color.rgb = COLOR_MUTED_GREEN
    arrow3.line.fill.background()

    tx_ret = slide.shapes.add_textbox(Inches(4.15), Inches(6.0), Inches(5.0), Inches(0.8))
    p_ret = tx_ret.text_frame.paragraphs[0]
    p_ret.text = "Real-time results & treatment returned to Mobile screen in < 2 seconds"
    p_ret.font.size = Pt(14)
    p_ret.font.bold = True
    p_ret.font.color.rgb = COLOR_PRIMARY_GREEN
    p_ret.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 7: Knowledge Base (Treatments)
    # ==========================================
    slide = prs.slides.add_slide(blank_slide_layout)
    add_header(slide, "Integrated Treatment Knowledge Base")

    pests = [
        ("Brown Planthopper", "Avoid excessive nitrogen fertilizers. Introduce natural enemies. Apply specific selective insecticides like Buprofezin."),
        ("Stem Borer", "Deploy pheromone traps. Maintain proper field drainage. Apply chemical controls like Carbofuran if threshold is exceeded."),
        ("Rice Bug", "Keep fields clear of host weeds. Hand-pick bugs or apply contact insecticides (e.g. Lambda-cyhalothrin) during early mornings."),
        ("Green Leafhopper", "Implement light traps to capture adults. Apply systemic chemical treatments to prevent secondary Tungro virus infection.")
    ]

    card_w = Inches(5.6)
    card_h = Inches(2.1)
    g_x = Inches(0.5)
    g_y = Inches(0.3)
    s_x = Inches(0.75)
    s_y = Inches(1.8)

    for i, (pest, treatment) in enumerate(pests):
        row = i // 2
        col = i % 2
        
        c_x = s_x + col * (card_w + g_x)
        c_y = s_y + row * (card_h + g_y)
        
        # Soft background card
        add_card(slide, c_x, c_y, card_w, card_h, COLOR_SOFT_GRAY, border_color=COLOR_MINT_LIGHT)
        
        tx = slide.shapes.add_textbox(c_x + Inches(0.15), c_y + Inches(0.15), card_w - Inches(0.3), card_h - Inches(0.3))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p_name = tf.paragraphs[0]
        p_name.text = "🎯  " + pest
        p_name.font.name = FONT_TITLE
        p_name.font.size = Pt(18)
        p_name.font.bold = True
        p_name.font.color.rgb = COLOR_PRIMARY_GREEN
        p_name.space_after = Pt(5)
        
        p_treat = tf.add_paragraph()
        p_treat.text = treatment
        p_treat.font.name = FONT_BODY
        p_treat.font.size = Pt(13)
        p_treat.font.color.rgb = COLOR_DARK_TEXT

    # ==========================================
    # SLIDE 8: Technology Stack Summary
    # ==========================================
    slide = prs.slides.add_slide(blank_slide_layout)
    add_header(slide, "The Complete Technology Stack")

    tech_cards = [
        ("Android Client App", "• Android SDK (Java)\n• OkHttp 3 HTTP Library\n• Jetpack Activity Results\n• FileProvider API", COLOR_MINT_LIGHT, COLOR_PRIMARY_GREEN),
        ("Python Backend", "• Flask Microframework\n• Werkzeug Secure Files\n• Python 3.10 Runtime\n• RESTful JSON API", COLOR_SOFT_GRAY, COLOR_DARK_TEXT),
        ("AI & Machine Learning", "• Ultralytics YOLOv8n\n• PyTorch Library\n• OpenCV for Image Handling\n• Custom Trained Weights", COLOR_SOFT_GRAY, COLOR_DARK_TEXT),
        ("Data & Infrastructure", "• Roboflow Datasets\n• Annotated Cereal Pests\n• Custom YAML Configuration\n• High-performance CUDA", COLOR_MINT_LIGHT, COLOR_PRIMARY_GREEN)
    ]

    c_w = Inches(5.6)
    c_h = Inches(2.1)
    g_x = Inches(0.5)
    g_y = Inches(0.3)
    s_x = Inches(0.75)
    s_y = Inches(1.8)

    for i, (title, items, bg, text_col) in enumerate(tech_cards):
        row = i // 2
        col = i % 2
        
        c_x = s_x + col * (c_w + g_x)
        c_y = s_y + row * (c_h + g_y)
        
        add_card(slide, c_x, c_y, c_w, c_h, bg)
        
        tx = slide.shapes.add_textbox(c_x + Inches(0.2), c_y + Inches(0.2), c_w - Inches(0.4), c_h - Inches(0.4))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = FONT_TITLE
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY_GREEN if bg == COLOR_MINT_LIGHT else COLOR_ACCENT_GOLD
        p_title.space_after = Pt(8)
        
        p_items = tf.add_paragraph()
        p_items.text = items
        p_items.font.name = FONT_BODY
        p_items.font.size = Pt(13)
        p_items.font.color.rgb = COLOR_DARK_TEXT

    # ==========================================
    # SLIDE 9: System Value & Impact
    # ==========================================
    slide = prs.slides.add_slide(blank_slide_layout)
    add_header(slide, "System Value & Real-World Impact")

    # Left side: key values
    txBox_val = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_val = txBox_val.text_frame
    tf_val.word_wrap = True

    values = [
        ("Universal Crop Protection", "Enables farmers to protect vital cereal grains with simple, free, and accessible mobile technology."),
        ("Environmental Sustainability", "Targeted recommendations avoid chemical saturation, preserving soils, beneficial insects, and water supplies."),
        ("Extensible Framework", "The model easily accommodates new pests (e.g. corn or potato diseases) by running additional transfer learning cycles.")
    ]

    for title, desc in values:
        p_t = tf_val.add_paragraph() if tf_val.paragraphs[0].text else tf_val.paragraphs[0]
        p_t.text = "• " + title
        p_t.font.name = FONT_TITLE
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_PRIMARY_GREEN
        p_t.space_after = Pt(4)

        p_d = tf_val.add_paragraph()
        p_d.text = desc
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = COLOR_DARK_TEXT
        p_d.space_after = Pt(20)

    # Right side: Summary Graphic Card
    add_card(slide, Inches(7.5), Inches(1.8), Inches(5.083), Inches(4.5), COLOR_PRIMARY_GREEN)
    txBox_val_g = slide.shapes.add_textbox(Inches(7.8), Inches(2.4), Inches(4.483), Inches(3.3))
    tf_vg = txBox_val_g.text_frame
    tf_vg.word_wrap = True

    val_lines = [
        ("DEMOCRATIZING", 18, COLOR_MINT_LIGHT, True),
        ("AGRICULTURAL SCIENCE", 24, COLOR_WHITE, True),
        ("\n---\n", 14, COLOR_ACCENT_GOLD, False),
        ("\nPutting custom object detection models and agricultural wisdom straight into the hands of local growers.", 16, COLOR_WHITE, False)
    ]

    for text, size, color, is_bold in val_lines:
        p = tf_vg.add_paragraph() if tf_vg.paragraphs[0].text else tf_vg.paragraphs[0]
        p.text = text
        p.font.name = FONT_TITLE
        p.font.size = Pt(size)
        p.font.bold = is_bold
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 10: Next Steps & Roadmap
    # ==========================================
    slide = prs.slides.add_slide(blank_slide_layout)
    add_header(slide, "Future Development Roadmap")

    # Left card: Next Steps
    add_card(slide, Inches(0.75), Inches(1.8), Inches(6.0), Inches(4.5), COLOR_SOFT_GRAY)
    tx_road_l = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.5), Inches(4.1))
    tf_rl = tx_road_l.text_frame
    tf_rl.word_wrap = True

    p_rl_title = tf_rl.paragraphs[0]
    p_rl_title.text = "Key Upgrades & Features"
    p_rl_title.font.name = FONT_TITLE
    p_rl_title.font.size = Pt(22)
    p_rl_title.font.bold = True
    p_rl_title.font.color.rgb = COLOR_PRIMARY_GREEN
    p_rl_title.space_after = Pt(15)

    roadmap_points = [
        ("Offline TFLite Inference:", "Converting the YOLOv8 model to TensorFlow Lite to execute on-device inference without internet."),
        ("Multi-Pest Overlays:", "Adding colored canvas boxes directly on the Android screen to showcase multiple pests in a single view."),
        ("Outbreak Heatmaps:", "Using mobile GPS logs to chart regional pest migration, enabling preemptive agricultural notices.")
    ]

    for lbl, text in roadmap_points:
        p = tf_rl.add_paragraph()
        p.text = f"🚀  {lbl} "
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK_TEXT
        
        run = p.add_run()
        run.text = text
        run.font.bold = False
        p.space_after = Pt(12)

    # Right Column: Big Closing Statement
    add_card(slide, Inches(7.25), Inches(1.8), Inches(5.333), Inches(4.5), COLOR_MINT_LIGHT)
    
    txBox_road_r = slide.shapes.add_textbox(Inches(7.5), Inches(2.5), Inches(4.833), Inches(3.3))
    tf_rr = txBox_road_r.text_frame
    tf_rr.word_wrap = True

    p_rr1 = tf_rr.paragraphs[0]
    p_rr1.text = "Empowering Farmers"
    p_rr1.font.name = FONT_TITLE
    p_rr1.font.size = Pt(28)
    p_rr1.font.bold = True
    p_rr1.font.color.rgb = COLOR_PRIMARY_GREEN
    p_rr1.alignment = PP_ALIGN.CENTER
    p_rr1.space_after = Pt(10)

    p_rr2 = tf_rr.add_paragraph()
    p_rr2.text = "Safeguarding Global Food Yields"
    p_rr2.font.name = FONT_TITLE
    p_rr2.font.size = Pt(20)
    p_rr2.font.bold = True
    p_rr2.font.color.rgb = COLOR_ACCENT_GOLD
    p_rr2.alignment = PP_ALIGN.CENTER
    p_rr2.space_after = Pt(20)

    p_rr3 = tf_rr.add_paragraph()
    p_rr3.text = "Implementing deep learning inside standard tools creates a resilient shield for global farming communities."
    p_rr3.font.name = FONT_BODY
    p_rr3.font.size = Pt(14)
    p_rr3.font.color.rgb = COLOR_DARK_TEXT
    p_rr3.alignment = PP_ALIGN.CENTER

    # Save Presentation
    output_filename = "Pest_Detection_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as '{output_filename}'!")

if __name__ == "__main__":
    create_presentation()
